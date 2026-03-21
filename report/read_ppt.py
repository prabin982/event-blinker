import sys
from pptx import Presentation

def extract_text(ppt_path):
    print(f"Reading {ppt_path}")
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"Error reading ppt: {e}")
        return
        
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            # print only first 100 chars to avoid spam, but we want all text actually
            print(f"Shape {j}: {repr(text)}")

if __name__ == '__main__':
    extract_text(sys.argv[1])
