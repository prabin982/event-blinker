import sys
from pptx import Presentation

def replace_text(shape, new_text):
    if not shape.has_text_frame:
        return
    # We want to clear the text frame but retain styling if possible
    # Just replacing text directly on the shape text frame might override some paragraph styles
    # We will replace text of the first run of the first paragraph, and clear the rest
    if len(shape.text_frame.paragraphs) > 0:
        p = shape.text_frame.paragraphs[0]
        if len(p.runs) > 0:
            p.runs[0].text = new_text
            for run_idx in range(1, len(p.runs)):
                p.runs[run_idx].text = ""
            for p_idx in range(1, len(shape.text_frame.paragraphs)):
                p = shape.text_frame.paragraphs[p_idx]
                for run in p.runs:
                    run.text = ""
        else:
            shape.text = new_text
    else:
        shape.text = new_text

def update_presentation(input_path, output_path):
    prs = Presentation(input_path)
    
    # Define replacements per slide and shape index
    # We map (slide_index, shape_index) -> new_text
    replacements = {
        # Slide 1
        (0, 7): "A REAL-TIME GEOSPATIAL EVENT AND RIDE-SHARING PLATFORM",
        
        # Slide 2 - Background
        (1, 9): " Information overload and fragmentation on social platforms.\n Difficulty in localized real-time event discovery (FOMO).\n Lack of integrated platforms addressing \"last-mile\" transportation.\n Event Blinker solves this by integrating maps, AI, and ride-sharing.",
        
        # Slide 3 - Problem
        (2, 10): "Information Staleness: Existing platforms lack real-time status updates.",
        (2, 14): "Geospatial Inefficiency: Hard to discover nearby events without intelligent map support.",
        (2, 18): "Transportation Barriers: Lack of localized integrated transport.",
        (2, 22): "Static Assistance: No context-aware, localized query resolution.",
        
        # Slide 5 - Literature
        (4, 0): "MEETUP: FOCUSED ON GROUPS; LACKS REAL-TIME UPDATES & RIDE INTEGRATION.\n\nEVENTBRITE: TICKET-ORIENTED; LIMITED LIVE STATUS OR REAL-TIME INTERACTION.\n\nFACEBOOK EVENTS: UNSTRUCTURED DATA; LACKS SPATIAL VERIFICATION AND DISCOVERY.",
        
        # Slide 6 - Tech Stack
        (5, 2): "FRONTEND: REACT NATIVE (EXPO), REACT / VITE",
        (5, 3): "BACKEND: NODE.JS, EXPRESS (11 ROUTE MODULES)",
        (5, 4): "DATABASE: POSTGRESQL + POSTGIS EXTENSION",
        (5, 5): "REAL-TIME & AI: SOCKET.IO, LLAMA 3.3 (OPENROUTER)",
        
        # Slide 7 - Scope
        (6, 6): "MOBILE APP FOR END USERS",
        (6, 7): "WEB PORTALS FOR ORGANIZERS & ADMIN",
        (6, 8): "AI VIRTUAL CONCIERGE (BLINKER AI)",
        (6, 9): "P2P RIDE-SHARING SYSTEM",
        (6, 10): "GEOSPATIAL EVENT ENGINE",
        
        # Slide 8 - Related Theory
        (7, 8): "Core Technologies",
        (7, 9): "PostGIS Geospatial Engine",
        (7, 17): "Three-Tier Client-Server Architecture",
        (7, 18): "Socket.io Bi-Directional Streaming",
        (7, 19): "Multi-Model AI Fallback Logic",
        
        # Slide 9 - Discovery
        (8, 8): "RADIUS-BASED FILTERING (ST_DWITHIN)",
        (8, 9): "INTERACTIVE MAPBOX GL DASHBOARD",
        (8, 10): "GEOSPATIAL EVENT SEARCH",
        (8, 11): "REAL-TIME MARKER UPDATES",
        (8, 12): "GEOSPATIAL DISCOVERY MODULE",
        
        # Slide 11 - Chat
        (10, 7): "AI CONCIERGE (LLAMA 3.3) FOR QUERIES",
        (10, 8): "INSTANT BROADCASTING VIA SOCKET.IO",
        (10, 9): "ATTENDEE AND ORGANIZER INTERACTION",
        (10, 10): "EVENT-SPECIFIC DISCUSSIONS",
        (10, 11): "LIVE CHAT & AI INTEGRATION",
        
        # Slide 14 - Ride Sharing
        (13, 7): "DYNAMIC PRICING & HAVERSINE METRICS",
        (13, 8): "3-STEP DRIVER VERIFICATION PROCESS",
        (13, 9): "P2P CUSTOM PRICE NEGOTIATIONS",
        (13, 10): "UNIFIED RIDE SHARING MODULE",
        (13, 11): "LIVE RIDE STATUS NOTIFICATIONS",
        
        # Slide 16 - Distance Calculation
        (15, 1): "DISTANCE IS REQUIRED FOR SPATIAL FILTERING & ACCURATE FARE LOGIC\n We utilize GPS coordinates (Latitude & Longitude).\n Haversine formula determines the shortest path over Earth's sphere.\n Output yields precise geographical distances in kilometers.",
        (15, 2): "INPUT DATA USED:\n Pickup coordinates (Lat, Lng)\n Drop-off coordinates (Lat, Lng)",
        
        # Slide 17 - Distance Calculation details
        (16, 0): "HOW DISTANCE AND FARE ARE CALCULATED\n\n Coordinates stored securely via PostGIS geometry types.\n PostGIS (ST_Distance) evaluates points for rapid querying.\n Haversine algorithm serves as backend standard for precise ride metrics.",
        (16, 1): "TOTAL FARE = BASE FARE (50) + (DISTANCE / KM x VEHICLE RATE)",
        
        # Slide 18 - Security
        (17, 7): "ADMIN VERIFICATION FOR ACTORS",
        (17, 8): "DRIVER LICENSE & DOCUMENT REVIEW",
        (17, 9): "BCRYPT HASHING & RATE LIMITING",
        (17, 10): "STATELESS JWT AUTHENTICATION",
        (17, 11): "PLATFORM SECURITY & INTEGRITY",
        
        # Slide 19 - Methodology -> Work Completed
        (18, 8): "Work Completed",
        (18, 9): "DATABASE & GEOSPATIAL ENGINE",
        (18, 11): "RIDE SHARING LIFESTYLE",
        (18, 15): "AI MULTI-MODEL INTEGRATION",
        (18, 16): "MOBLE, WEB & ADMIN PORTALS",
        (18, 17): "PRODUCTION DEPLOYMENT (RENDER)",
        (18, 19): "VARIOUS ROUTE MODULES COMPLETED",
        
        # Slide 20 - Project Progress
        (19, 1): "PROJECT PROGRESS (ACHIEVEMENTS)\n CORE BACKEND ARCHITECTURE (11 MODULES)\n POSTGIS GEOSPATIAL EVENT DISCOVERY\n BLINKER AI CONCIERGE ACTIVE\n P2P RIDE-SHARING FULLY IMPLEMENTED\n CROSS-PLATFORM SYSTEM DEPLOYED",
        
        # Slide 21 - Challenges
        (20, 1): " POSTGIS GIST INDEX OPTIMIZATION\n SOCKET.IO CORS COMPATIBILITY\n AI HALLUCINATION & CONTEXT GROUNDING\n MULTI-MODEL FALLBACK RELIABILITY\n SCHEMA EVOLUTION MANAGEMENT",
        (20, 2): "CHALLENGES RESOLVED",
        
        # Slide 22 - Conclusion
        (21, 0): "EVENT BLINKER SEAMLESSLY MERGES EVENT DISCOVERY, LOCALIZED AI ASSISTANCE, AND LOGISTICS INTO ONE UNIFIED ECOSYSTEM. UTILIZING POSTGIS AND SOCKET.IO ENSURES HIGH-PERFORMANCE, REAL-TIME EXPERIENCES WHILE CATERING TO LOCAL TRANSPORT NEEDS.",
        (21, 1): "CONCLUSION",
    }
    
    for (slide_idx, shape_idx), text in replacements.items():
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            if shape_idx < len(slide.shapes):
                shape = slide.shapes[shape_idx]
                replace_text(shape, text)
            else:
                print(f"Warning: Slide {slide_idx} does not have shape {shape_idx}")
        else:
            print(f"Warning: Presentation does not have slide {slide_idx}")
            
    # Try to insert objective values into Slide 4 (index 3) if there is an empty text box we can use.
    # From read_ppt output, slide 4 (index 3) has shapes:
    # Shape 7: 'Introduction'
    # Shape 8: 'Objectives'
    # Shape 13: '3'
    # Let's add text to shape 0, 1, 2, or 4 if they are empty
    slide_3 = prs.slides[3]
    obj_text = " Real-time event discovery with Mapbox GL & PostGIS.\n AI Virtual Concierge (Blinker AI) via LLAMA 3.3.\n P2P ride-sharing with Haversine distance & NPR dynamic pricing.\n Verification portals for Events, Organizers, and Riders."
    for idx in [0, 1, 2, 4]:
        if idx < len(slide_3.shapes) and slide_3.shapes[idx].has_text_frame:
            replace_text(slide_3.shapes[idx], obj_text)
            break
            
    prs.save(output_path)
    print(f"Saved modified presentation to {output_path}")

if __name__ == '__main__':
    input_ppt = sys.argv[1]
    output_ppt = sys.argv[2]
    update_presentation(input_ppt, output_ppt)
